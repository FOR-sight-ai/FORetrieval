from __future__ import annotations

import os
import re
import json
import mimetypes
from datetime import datetime, timezone

from pathlib import Path
from typing import Any, Callable, Dict, Optional, List, Union

from pydantic import BaseModel, Field
import asyncio
import inspect

# PydanticAI – cœur
from pydantic_ai import Agent

# OpenAI / OpenRouter / Ollama (API OpenAI-compatible)
from pydantic_ai.models.openai import OpenAIModel
from pydantic_ai.providers.openai import OpenAIProvider

# Mistral (client natif PydanticAI)
try:
    from pydantic_ai.models.mistral import MistralModel
    from pydantic_ai.providers.mistral import MistralProvider
    _MISTRAL_AVAILABLE = True
except Exception:
    _MISTRAL_AVAILABLE = False

# Langue & parsers fichiers (optionnels)
try:
    from langdetect import detect as lang_detect, DetectorFactory
    DetectorFactory.seed = 0
except Exception:
    lang_detect = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
except Exception:
    Image = None


class AIMetadata(BaseModel):
    author: Optional[str] = Field(None)
    title: Optional[str] = Field(None)
    language: Optional[str] = Field(None)
    tags: List[str] = Field(default_factory=list)
    document_type: str
    short_description: str


def _iso_utc(ts: float) -> str:
    return datetime.fromtimestamp(ts, tz=timezone.utc).replace(microsecond=0).isoformat()


def _guess_mime(path: Path) -> str:
    mime, _ = mimetypes.guess_type(str(path))
    return mime or "application/octet-stream"


def _pdf_info(path: Path) -> Dict[str, Union[int, str, None]]:
    out: Dict[str, Union[int, str, None]] = {"page_count": None, "author": None, "title": None, "text_preview": ""}
    if not PdfReader:
        return out
    try:
        r = PdfReader(str(path))
        out["page_count"] = len(r.pages)
        meta = (r.metadata or {})
        out["author"] = getattr(meta, "author", None) or meta.get("/Author")
        out["title"]  = getattr(meta, "title", None)  or meta.get("/Title")
        chunks = []
        for i in range(min(3, len(r.pages))):
            try:
                chunks.append(r.pages[i].extract_text() or "")
            except Exception:
                pass
        out["text_preview"] = "\n".join(chunks).strip()
    except Exception:
        pass
    return out


def _docx_text(path: Path) -> str:
    if not docx:
        return ""
    try:
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs if p.text).strip()
    except Exception:
        return ""


def _pptx_text(path: Path) -> str:
    if not Presentation:
        return ""
    try:
        prs = Presentation(str(path))
        t = []
        for s in prs.slides:
            for shp in s.shapes:
                if hasattr(shp, "text") and shp.text:
                    t.append(shp.text)
        return "\n".join(t).strip()
    except Exception:
        return ""


def _image_dims(path: Path):
    if not Image:
        return (None, None)
    try:
        with Image.open(path) as im:
            return im.size
    except Exception:
        return (None, None)


def _text_preview(path: Path, ext: str) -> str:
    ext = ext.lower()
    if ext == ".pdf":
        return _pdf_info(path)["text_preview"] or ""
    if ext in {".docx"}:
        return _docx_text(path)
    if ext in {".pptx", ".ppt"}:
        return _pptx_text(path)
    if ext in {".txt", ".md", ".csv", ".json", ".yaml", ".yml"}:
        try:
            return Path(path).read_text(encoding="utf-8", errors="ignore")[:10000]
        except Exception:
            return ""
    return ""


def _fallback_lang(text: str) -> Optional[str]:
    if not text or not text.strip() or not lang_detect:
        return None
    try:
        return lang_detect(text[:4000])
    except Exception:
        return None


AI_SYSTEM_PROMPT = (
    "You are a meticulous document analyst. "
    "Given raw text preview and basic file hints, you MUST output a compact, factual summary; "
    "infer document_type from common categories; infer language; "
    "return short, meaningful tags (3–8 max). "
    "Keep author/title null if unknown; avoid hallucinating."
)

AI_INSTRUCTIONS = (
    "Analyse the document preview and hints to produce strictly the fields of AIMetadata. "
    "Language should be a short code like 'fr' or 'en'. "
    "Document type: choose a concise label (e.g., 'technical note', 'scientific paper', 'accounting', 'contract', "
    "'manual', 'report', 'invoice', 'specification', 'slide deck'). "
    "Short description: 1–3 sentences. Keep it helpful and neutral. "
    "Tags: 3–8 short topic tags."
)


def _build_model_from_cfg(ai_cfg: Dict[str, Any]):
    """
    Construit dynamiquement un modèle PydanticAI selon le provider choisi.

    ai_cfg attend par ex.:
      {
        "provider": "openai" | "openrouter" | "mistral" | "mistralai" | "ollama",
        "api_key": "...",
        "base_url": "...",     # optionnel, défauts fournis ci-dessous
        "name": "model-id"     # optionnel, défauts fournis ci-dessous
      }
    """
    provider_raw = (ai_cfg.get("provider") or "openai").lower()
    name = ai_cfg.get("name")
    base_url = ai_cfg.get("base_url")
    api_key = ai_cfg.get("api_key")

    if provider_raw in ("mistral", "mistralai"):
        if not _MISTRAL_AVAILABLE:
            raise RuntimeError(
                "Le provider 'mistral' est demandé mais les dépendances Mistral pour pydantic-ai ne sont pas installées.\n"
                "Installe:  pip install \"pydantic-ai-slim[mistral]\"\n"
                "ou       : pip install \"pydantic-ai[mistral]\""
            )
        # Defaults Mistral
        base_url = base_url or "https://api.mistral.ai/v1"
        api_key = api_key or os.getenv("MISTRAL_API_KEY")
        if not name:
            # Choix par défaut raisonnable
            name = "mistral-small-latest"

        provider = MistralProvider(
            api_key=api_key,
            base_url=base_url,
        )
        return MistralModel(name, provider=provider)

    elif provider_raw == "openrouter":
        base_url = base_url or "https://openrouter.ai/api/v1"
        api_key = api_key or os.getenv("OPENROUTER_API_KEY")
        if not name:
            # Important: namespace requis chez OpenRouter
            name = "mistralai/mistral-small-latest"
        provider = OpenAIProvider(api_key=api_key, base_url=base_url)
        return OpenAIModel(model_name=name, provider=provider)

    elif provider_raw == "ollama":
        # Ollama expose une API OpenAI-compatible en local
        base_url = base_url or "http://localhost:11434/v1"
        api_key = api_key or os.getenv("OLLAMA_API_KEY") or "ollama"
        if not name:
            name = "llama3.1"  # ajuster selon vos tags locaux
        provider = OpenAIProvider(api_key=api_key, base_url=base_url)
        return OpenAIModel(model_name=name, provider=provider)

    elif provider_raw in ("openai", "openaiapi", "openai-api"):
        base_url = base_url or os.getenv("OPENAI_BASE_URL")  # généralement None => api.openai.com
        api_key = api_key or os.getenv("OPENAI_API_KEY")
        if not name:
            name = "gpt-4o-mini"
        provider = OpenAIProvider(api_key=api_key, base_url=base_url)
        return OpenAIModel(model_name=name, provider=provider)

    else:
        raise ValueError(f"Provider inconnu: {provider_raw}")


def ai_metadata_provider_factory(ai_cfg: Optional[Dict[str, Any]]) -> Callable[[Path], Dict[str, Any]]:
    """
    Retourne un callable(Path) -> Dict métadonnées complètes.
    - Sans ai_cfg => provider 'no-LLM' (remplit seulement les champs bruts).
    - Avec ai_cfg => construit un Agent PydanticAI avec le bon modèle selon le provider,
                     instructions=AI_SYSTEM_PROMPT, output_type=AIMetadata.
    """

    # ----- Sans IA : juste les champs "bruts" -----
    if not ai_cfg:
        def _no_ai_provider(p: Path) -> Dict[str, Any]:
            p = p.resolve()
            stat = p.stat()
            ext = p.suffix.lower()
            mime = _guess_mime(p)
            md: Dict[str, Any] = {
                "source_path": str(p),
                "stem": p.stem,
                "ext": ext,
                "mime": mime,
                "mtime": _iso_utc(stat.st_mtime),
                "page_count": None,
                "image_width": None,
                "image_height": None,
                "author": None,
                "title": None,
                "language": None,
                "tags": [],
                "document_type": "unknown",
                "short_description": "",
            }
            if ext == ".pdf":
                pdf_meta = _pdf_info(p)
                md["page_count"] = pdf_meta.get("page_count")
                if pdf_meta.get("author"):
                    md["author"] = str(pdf_meta["author"])
                if pdf_meta.get("title"):
                    md["title"] = str(pdf_meta["title"])
            if ext in {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}:
                w, h = _image_dims(p)
                md["image_width"], md["image_height"] = w, h
            return md
        return _no_ai_provider

    # ----- Avec IA : construit le modèle & l’agent PydanticAI -----
    model = _build_model_from_cfg(ai_cfg)

    # NB: structure PydanticAI canonique : instructions = système, output_type = schéma Pydantic
    # Ajout d’un param "tools": [] pour éviter toute tentative de tool-calling si le backend ne supporte pas
    agent = Agent(
        model=model,
        instructions=AI_SYSTEM_PROMPT,
        output_type=AIMetadata,
        model_settings={'temperature': 0, 'tools': []}
    )

    def _provider(p: Path) -> Dict[str, Any]:
        p = p.resolve()
        stat = p.stat()
        ext = p.suffix.lower()
        mime = _guess_mime(p)

        md: Dict[str, Any] = {
            "source_path": str(p),
            "stem": p.stem,
            "ext": ext,
            "mime": mime,
            "mtime": _iso_utc(stat.st_mtime),
            "page_count": None,
            "image_width": None,
            "image_height": None,
            "author": None,
            "title": None,
            "language": None,
            "tags": [],
            "document_type": "unknown",
            "short_description": "",
        }

        if ext == ".pdf":
            pdf_meta = _pdf_info(p)
            md["page_count"] = pdf_meta.get("page_count")
            if pdf_meta.get("author"):
                md["author"] = str(pdf_meta["author"])
            if pdf_meta.get("title"):
                md["title"] = str(pdf_meta["title"])

        if ext in {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"}:
            w, h = _image_dims(p)
            md["image_width"], md["image_height"] = w, h

        preview = _text_preview(p, ext)
        lang_hint = _fallback_lang(preview) if preview else None

        # Contexte utilisateur (passé dans le prompt de run)
        ctx = {
            "file_name": p.name,
            "mime": mime,
            "ext": ext,
            "pdf_author_hint": md.get("author"),
            "pdf_title_hint": md.get("title"),
            "detected_lang_hint": lang_hint,
            "text_preview": (preview[:12000] if preview else ""),
        }
        user_prompt = (
            AI_INSTRUCTIONS
            + "\n\n### CONTEXT (JSON)\n"
            + json.dumps(ctx, ensure_ascii=False)
            + "\n\n### OUTPUT FORMAT\n"
            + "Respond strictly with the required fields; do not invent author/title if unknown."
        )

        # ✅ API PydanticAI canonique : run(prompt) -> result.output (typé output_type)
        try:
            result = agent.run_sync(user_prompt)  # pydantic-ai >= 0.0.15 (selon versions)
        except AttributeError:
            result = asyncio.run(agent.run_sync(user_prompt))  # fallback universel

        ai: AIMetadata = result.output

        # Fusion des champs
        if ai.author:
            md["author"] = ai.author
        if ai.title:
            md["title"] = ai.title
        md["language"] = ai.language or lang_hint
        md["tags"] = ai.tags or []
        md["document_type"] = ai.document_type
        md["short_description"] = ai.short_description
        print("result", md)
        return md

    return _provider
